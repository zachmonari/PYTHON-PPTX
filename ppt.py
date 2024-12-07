# Re-initialize the process to create the PowerPoint presentation after code execution reset

from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Comprehensive Analysis and Design of Structural Elements in Mechanical Engineering"
subtitle.text = (
    "A Detailed Study of Beams, Struts, Plates, and Shells\n"
    "Presented by: Zachary Monari, Chacha Francis, Kelvin Agwata, Isaac Wanje, Norman Sande, "
    "Ian Otieno, Daniel Ndegwa\n"
    "Dedan Kimathi University of Technology\n"
    "BSc Mechanical Engineering"
)

# Slide 2: Abstract
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Abstract"
content = slide.placeholders[1]
content.text = (
    "This term paper presents a comprehensive analysis and design of structural elements "
    "in mechanical engineering, focusing on beams, struts, plates, and shells. It evaluates "
    "torsional shear stresses, stability in columns, and stress distributions in plates and shells. "
    "Analytical results are compared with empirical codes to validate design accuracy."
)

# Slide 3: Objectives
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Objectives"
content = slide.placeholders[1]
content.text = (
    "- Evaluate shear stresses, deflection, and torsion in thin-walled structures.\n"
    "- Design and analyze struts under various conditions.\n"
    "- Solve statically indeterminate systems analytically.\n"
    "- Design plates and shells using plate and shell mechanics."
)

# Slide 4: Theoretical Background (Part 1)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Theoretical Background: Part 1"
content = slide.placeholders[1]
content.text = (
    "- Shear stresses in non-circular sections exhibit non-uniform distributions.\n"
    "- Warping deformation occurs in thin-walled sections.\n"
    "Key Formula:\n"
    "   τ = T / (2Am * t)"
)

# Slide 5: Theoretical Background (Part 2)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Theoretical Background: Part 2"
content = slide.placeholders[1]
content.text = (
    "- Columns fail due to buckling or crushing, depending on their geometry.\n"
    "- Euler’s Formula for critical load:\n"
    "   Pcr = π²EI / (KL)²\n"
    "- Empirical codes improve practical designs."
)

# Slide 6: Theoretical Background (Part 3)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Theoretical Background: Part 3"
content = slide.placeholders[1]
content.text = (
    "- Plates experience bending stresses; shells exhibit membrane stresses.\n"
    "Key Formulas:\n"
    "   Plate Deflection: D∇⁴w = q\n"
    "   Shell Membrane Stress: σm = pR / 2t"
)

# Slide 7: Discussion (Part 2)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Discussion: Struts and Columns"
content = slide.placeholders[1]
content.text = (
    "- Buckling sensitivity to end constraints and imperfections.\n"
    "- Eccentric loading amplifies stress; secant formula provides corrections.\n"
    "- Pre-stressed columns and lateral bracing enhance stability."
)

# Slide 8: Discussion (Part 3)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Discussion: Plates and Shells"
content = slide.placeholders[1]
content.text = (
    "- Plates experience maximum bending stress at the center.\n"
    "- Shell curvature reduces bending stresses, ensuring efficient load distribution.\n"
    "- Honeycomb core plates improve stiffness and reduce weight."
)

# Slide 9: Challenges
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Challenges"
content = slide.placeholders[1]
content.text = (
    "- Stress distribution in non-circular sections requires numerical methods.\n"
    "- Column imperfections significantly affect buckling loads.\n"
    "- Balancing weight and strength in thin plates and shells."
)

# Slide 10: Proposed Solutions
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Proposed Solutions"
content = slide.placeholders[1]
content.text = (
    "- Use FEA tools for precise analysis (e.g., ANSYS).\n"
    "- Employ lightweight, high-strength materials like composites.\n"
    "- Follow empirical codes (Eurocode 3, ASME Section VIII) for reliability."
)

# Slide 11: Conclusion
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1]
content.text = (
    "- Non-circular sections require reinforcement to handle shear stress.\n"
    "- Stability in columns depends on mitigating imperfections.\n"
    "- Plates and shells benefit from optimized curvature and advanced materials."
)

# Slide 12: References
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "References"
content = slide.placeholders[1]
content.text = (
    "1. Beer, F. P., et al. (2017). Mechanics of Materials.\n"
    "2. Reddy, J. N. (2007). Theory and Analysis of Elastic Plates and Shells.\n"
    "3. Timoshenko, S. P., & Woinowsky-Krieger, S. (1959). Theory of Plates and Shells.\n"
    "4. Ugural, A. C., & Fenster, S. K. (2003). Advanced Strength and Applied Elasticity."
)

# Save the presentation to a file
output_path = "Structural_Elements_Term_Paper.pptx"
presentation.save(output_path)
#output_path
